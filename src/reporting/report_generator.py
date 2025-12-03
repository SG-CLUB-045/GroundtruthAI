"""
Report Generation Module
Creates PDF and PowerPoint reports with charts, tables, and insights
"""

import os
import logging
import re
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    pass

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image as RLImage
    from reportlab.lib import colors
    from reportlab.lib.colors import HexColor
    def ReportLabRGB(r, g, b):
        return HexColor(f"#{r:02x}{g:02x}{b:02x}")
    REPORTLAB_AVAILABLE = True
except ImportError:
    ReportLabRGB = None
    REPORTLAB_AVAILABLE = False

import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
import pandas as pd

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def clean_text_for_presentation(text: str) -> str:
    """
    Clean text for presentation by removing markdown and formatting characters.
    
    Args:
        text: Raw text that may contain markdown
        
    Returns:
        Cleaned text suitable for presentation
    """
    if not text:
        return ""
    
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    text = re.sub(r'_([^_]+)_', r'\1', text)
    text = re.sub(r'\*+', '', text)
    text = re.sub(r' +', ' ', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    lines = [line.strip() for line in text.split('\n')]
    text = '\n'.join(lines)
    
    return text.strip()


class ReportGenerator:
    """
    Generates PDF and PowerPoint reports with data visualizations and insights.
    """

    def __init__(self, output_dir: str = "./output"):
        """
        Initialize ReportGenerator.
        
        Args:
            output_dir: Directory for output files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.charts: List[str] = []

    def create_chart(
        self,
        data: Dict[str, Any],
        chart_type: str = "bar",
        title: str = "Chart",
        xlabel: str = "",
        ylabel: str = "",
        figsize: Tuple[int, int] = (10, 6)
    ) -> str:
        """
        Create a chart and save it as an image.
        
        Args:
            data: Dictionary with chart data
            chart_type: Type of chart ('bar', 'line', 'pie', 'scatter')
            title: Chart title (will be updated with actual item count)
            xlabel: X-axis label
            ylabel: Y-axis label
            figsize: Figure size (width, height)
            
        Returns:
            Path to saved chart image
        """
        try:
            fig, ax = plt.subplots(figsize=figsize)
            fig.patch.set_facecolor('white')
            
            if chart_type == "bar":
                total_items = len(data)
                x_data = list(data.keys())[:20]
                y_data = list(data.values())[:20]
                actual_count = len(x_data)
                
                if "Top" in title and total_items != actual_count:
                    title = title.replace("Top 10", f"Top {actual_count}").replace("Top 20", f"Top {actual_count}")
                elif actual_count < total_items and "Top" not in title:
                    title = f"{title} ({actual_count} items shown)"
                
                ax.bar(x_data, y_data, color='#2E86AB', alpha=0.8)
                ax.set_xticks(range(len(x_data)))
                ax.set_xticklabels(x_data, rotation=45, ha='right')
            
            elif chart_type == "line":
                x_data = list(data.keys())
                y_data = list(data.values())
                ax.plot(range(len(x_data)), y_data, marker='o', color='#2E86AB', linewidth=2)
                ax.set_xticks(range(len(x_data)))
                ax.set_xticklabels(x_data, rotation=45, ha='right')
            
            elif chart_type == "pie":
                labels = list(data.keys())[:10]
                sizes = list(data.values())[:10]
                colors_palette = ['#2E86AB', '#A23B72', '#F18F01', '#C73E1D', '#6A994E',
                                 '#BC4749', '#2C7873', '#4C9A9B', '#8B0000', '#FF6B6B']
                ax.pie(sizes, labels=labels, autopct='%1.1f%%', colors=colors_palette, startangle=90)
            
            elif chart_type == "scatter":
                x_data = list(range(len(data)))
                y_data = list(data.values())
                ax.scatter(x_data, y_data, s=100, alpha=0.6, color='#2E86AB')
            
            ax.set_title(title, fontsize=14, fontweight='bold', pad=20)
            if xlabel:
                ax.set_xlabel(xlabel, fontsize=11)
            if ylabel:
                ax.set_ylabel(ylabel, fontsize=11)
            
            ax.grid(True, alpha=0.3)
            plt.tight_layout()
            
            chart_path = self.output_dir / f"chart_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.png"
            plt.savefig(str(chart_path), dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            self.charts.append(str(chart_path))
            logger.info(f"Created {chart_type} chart: {chart_path}")
            return str(chart_path)
            
        except Exception as e:
            logger.error(f"Error creating chart: {str(e)}")
            return ""

    def create_table_image(
        self,
        data: Dict[str, List[Any]],
        title: str = "Data Table",
        rows_limit: int = 10
    ) -> str:
        """
        Create a table as an image.
        
        Args:
            data: Dictionary with column names as keys and lists as values
            title: Table title
            rows_limit: Maximum rows to display
            
        Returns:
            Path to saved table image
        """
        try:
            # Create DataFrame from data
            df = pd.DataFrame(data)
            df = df.head(rows_limit)
            
            fig, ax = plt.subplots(figsize=(12, max(4, len(df) * 0.3)))
            ax.axis('tight')
            ax.axis('off')
            
            # Create table
            table = ax.table(
                cellText=df.values,
                colLabels=df.columns,
                cellLoc='center',
                loc='center',
                colWidths=[1 / len(df.columns) for _ in df.columns]
            )
            
            table.auto_set_font_size(False)
            table.set_fontsize(9)
            table.scale(1, 1.5)
            
            for i in range(len(df.columns)):
                table[(0, i)].set_facecolor('#2E86AB')
                table[(0, i)].set_text_props(weight='bold', color='white')
            
            for i in range(1, len(df) + 1):
                for j in range(len(df.columns)):
                    if i % 2 == 0:
                        table[(i, j)].set_facecolor('#F0F0F0')
                    else:
                        table[(i, j)].set_facecolor('#FFFFFF')
            
            plt.title(title, fontsize=12, fontweight='bold', pad=20)
            plt.tight_layout()
            
            table_path = self.output_dir / f"table_{datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.png"
            plt.savefig(str(table_path), dpi=150, bbox_inches='tight', facecolor='white')
            plt.close()
            
            logger.info(f"Created table image: {table_path}")
            return str(table_path)
            
        except Exception as e:
            logger.error(f"Error creating table: {str(e)}")
            return ""

    def generate_pdf_report(
        self,
        title: str,
        subtitle: str,
        sections: Dict[str, str],
        tables: Optional[List[Dict[str, Any]]] = None,
        charts: Optional[List[str]] = None,
        insights: Optional[Dict[str, Any]] = None,
        output_filename: str = "report.pdf"
    ) -> str:
        """
        Generate a comprehensive PDF business report.
        
        Args:
            title: Report title
            subtitle: Report subtitle
            sections: Dictionary of section_name: section_content
            tables: List of table data dictionaries
            charts: List of chart image paths
            insights: Dictionary of insights
            output_filename: Output PDF filename
            
        Returns:
            Path to generated PDF
        """
        try:
            if not REPORTLAB_AVAILABLE or ReportLabRGB is None:
                raise ImportError("ReportLab not properly imported")
            
            output_path = self.output_dir / output_filename
            doc = SimpleDocTemplate(str(output_path), pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=28,
                textColor=ReportLabRGB(46, 134, 171),
                spaceAfter=12,
                alignment=1,
                fontName='Helvetica-Bold'
            )
            
            subtitle_style = ParagraphStyle(
                'CustomSubtitle',
                parent=styles['Normal'],
                fontSize=16,
                textColor=ReportLabRGB(100, 100, 100),
                spaceAfter=20,
                alignment=1,
                fontName='Helvetica'
            )
            
            heading2_style = ParagraphStyle(
                'CustomHeading2',
                parent=styles['Heading2'],
                fontSize=18,
                textColor=ReportLabRGB(46, 134, 171),
                spaceAfter=12,
                spaceBefore=12,
                fontName='Helvetica-Bold'
            )
            
            heading3_style = ParagraphStyle(
                'CustomHeading3',
                parent=styles['Heading3'],
                fontSize=14,
                textColor=ReportLabRGB(70, 70, 70),
                spaceAfter=8,
                spaceBefore=8,
                fontName='Helvetica-Bold'
            )
            
            normal_style = ParagraphStyle(
                'CustomNormal',
                parent=styles['Normal'],
                fontSize=11,
                textColor=ReportLabRGB(50, 50, 50),
                spaceAfter=6,
                leading=14,
                fontName='Helvetica'
            )
            
            # Cover Page
            story.append(Spacer(1, 1.5*inch))
            story.append(Paragraph(title, title_style))
            story.append(Spacer(1, 0.3*inch))
            story.append(Paragraph(subtitle, subtitle_style))
            story.append(Spacer(1, 0.5*inch))
            
            metadata = f"Generated on {datetime.now().strftime('%B %d, %Y at %H:%M:%S')}"
            story.append(Paragraph(metadata, normal_style))
            story.append(PageBreak())
            
            story.append(Paragraph("Table of Contents", heading2_style))
            toc_items = list(sections.keys())
            if insights:
                toc_items.extend(["Key Insights", "Key Findings", "Recommendations"])
            if tables:
                toc_items.append("Data Tables")
            if charts:
                toc_items.append("Visualizations")
            
            for i, item in enumerate(toc_items[:10], 1):  # Limit to 10 items
                story.append(Paragraph(f"{i}. {item}", normal_style))
            story.append(PageBreak())
            
            # Add sections with enhanced formatting
            for section_name, section_content in sections.items():
                story.append(Spacer(1, 0.2*inch))
                story.append(Paragraph(section_name, heading2_style))
                story.append(Spacer(1, 0.15*inch))
                
                section_content_clean = clean_text_for_presentation(str(section_content))
                content_lines = section_content_clean.split('\n')
                for line in content_lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                        clean_line = line.lstrip('•-* ').strip()
                        clean_line = clean_text_for_presentation(clean_line)
                        if clean_line:
                            story.append(Paragraph(f"• {clean_line}", normal_style))
                            story.append(Spacer(1, 0.08*inch))
                    elif line and line[0].isdigit() and ('.' in line[:3] or ')' in line[:3]):
                        clean_line = line.lstrip('0123456789.) ').strip()
                        if clean_line:
                            story.append(Paragraph(f"• {clean_line}", normal_style))
                            story.append(Spacer(1, 0.08*inch))
                    elif ':' in line and not line.startswith(' '):
                        parts = line.split(':', 1)
                        if len(parts) == 2:
                            key = parts[0].strip()
                            value = parts[1].strip()
                            story.append(Paragraph(f"<b>{key}:</b> {value}", normal_style))
                            story.append(Spacer(1, 0.08*inch))
                        else:
                            story.append(Paragraph(line, normal_style))
                            story.append(Spacer(1, 0.08*inch))
                    else:
                        if len(line) > 150:
                            sentences = line.split('. ')
                            for i, sentence in enumerate(sentences):
                                if sentence.strip():
                                    text = sentence.strip()
                                    if not text.endswith('.'):
                                        text += '.'
                                    story.append(Paragraph(text, normal_style))
                                    if i < len(sentences) - 1:
                                        story.append(Spacer(1, 0.05*inch))
                        else:
                            story.append(Paragraph(line, normal_style))
                            story.append(Spacer(1, 0.08*inch))
                
                story.append(Spacer(1, 0.3*inch))
            
            if insights and "Strategic Recommendations" in insights:
                story.append(PageBreak())
                story.append(Paragraph("Prioritized Strategic Recommendations", heading2_style))
                story.append(Spacer(1, 0.2*inch))
                
                strategic_recs = insights["Strategic Recommendations"]
                if isinstance(strategic_recs, list):
                    for idx, rec in enumerate(strategic_recs[:10], 1):
                        rec_clean = clean_text_for_presentation(str(rec))
                        if rec_clean:
                            story.append(Paragraph(f"{idx}. {rec_clean}", normal_style))
                            story.append(Spacer(1, 0.12*inch))
                story.append(Spacer(1, 0.2*inch))
            
            if insights:
                story.append(PageBreak())
                story.append(Paragraph("Strategic Insights & Analysis", heading2_style))
                story.append(Spacer(1, 0.2*inch))
                
                if "Executive Summary" in insights:
                    story.append(Paragraph("Executive Summary", heading3_style))
                    story.append(Spacer(1, 0.15*inch))
                    exec_summary = insights["Executive Summary"]
                    if isinstance(exec_summary, str):
                        exec_summary = clean_text_for_presentation(exec_summary)
                        summary_paragraphs = exec_summary.split('\n\n')
                        for para in summary_paragraphs:
                            para = para.strip()
                            if para:
                                if len(para) > 200:
                                    sentences = para.split('. ')
                                    for sentence in sentences:
                                        if sentence.strip():
                                            text = sentence.strip()
                                            if not text.endswith('.'):
                                                text += '.'
                                            story.append(Paragraph(text, normal_style))
                                            story.append(Spacer(1, 0.05*inch))
                                else:
                                    story.append(Paragraph(para, normal_style))
                                    story.append(Spacer(1, 0.1*inch))
                    story.append(Spacer(1, 0.2*inch))
                
                if "Key Findings" in insights:
                    story.append(Paragraph("Key Strategic Findings", heading3_style))
                    story.append(Spacer(1, 0.15*inch))
                    findings = insights["Key Findings"]
                    if isinstance(findings, list):
                        for idx, finding in enumerate(findings, 1):
                            finding_clean = clean_text_for_presentation(str(finding))
                            if finding_clean:
                                story.append(Paragraph(f"{idx}. {finding_clean}", normal_style))
                                story.append(Spacer(1, 0.12*inch))
                    else:
                        findings_text = clean_text_for_presentation(str(findings))
                        story.append(Paragraph(findings_text, normal_style))
                    story.append(Spacer(1, 0.2*inch))
                
                if "Recommendations" in insights:
                    story.append(Paragraph("AI-Generated Strategic Recommendations", heading3_style))
                    story.append(Spacer(1, 0.15*inch))
                    recommendations = insights["Recommendations"]
                    if isinstance(recommendations, list):
                        for idx, rec in enumerate(recommendations, 1):
                            rec_clean = clean_text_for_presentation(str(rec))
                            if rec_clean:
                                story.append(Paragraph(f"{idx}. {rec_clean}", normal_style))
                                story.append(Spacer(1, 0.12*inch))
                    else:
                        rec_text = clean_text_for_presentation(str(recommendations))
                        story.append(Paragraph(rec_text, normal_style))
                    story.append(Spacer(1, 0.2*inch))
            
            if tables:
                story.append(PageBreak())
                story.append(Paragraph("Data Analysis Tables", heading2_style))
                story.append(Spacer(1, 0.2*inch))
                
                for table_data in tables:
                    story.append(Paragraph(table_data.get('title', 'Data Table'), heading3_style))
                    
                    table_content = table_data.get('data', [])
                    if table_content and len(table_content) > 0:
                        num_cols = len(table_content[0])
                        col_widths = [6.5*inch / num_cols] * num_cols
                        
                        table_obj = Table(table_content, colWidths=col_widths)
                        table_style = TableStyle([
                            ('BACKGROUND', (0, 0), (-1, 0), ReportLabRGB(46, 134, 171)),
                            ('TEXTCOLOR', (0, 0), (-1, 0), ReportLabRGB(255, 255, 255)),
                            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                            ('FONTSIZE', (0, 0), (-1, 0), 11),
                            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                            ('TOPPADDING', (0, 0), (-1, 0), 12),
                            ('BACKGROUND', (0, 1), (-1, -1), ReportLabRGB(255, 255, 255)),
                            ('TEXTCOLOR', (0, 1), (-1, -1), ReportLabRGB(50, 50, 50)),
                            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                            ('FONTSIZE', (0, 1), (-1, -1), 10),
                            ('GRID', (0, 0), (-1, -1), 0.5, ReportLabRGB(200, 200, 200)),
                            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [ReportLabRGB(245, 245, 245), ReportLabRGB(255, 255, 255)]),
                            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ])
                        table_obj.setStyle(table_style)
                        story.append(table_obj)
                    story.append(Spacer(1, 0.3*inch))
            
            if charts:
                story.append(PageBreak())
                story.append(Paragraph("Data Visualizations", heading2_style))
                story.append(Spacer(1, 0.2*inch))
                
                for idx, chart_path in enumerate(charts, 1):
                    if os.path.exists(chart_path):
                        try:
                            story.append(Paragraph(f"Chart {idx}", heading3_style))
                            story.append(Spacer(1, 0.1*inch))
                            story.append(RLImage(chart_path, width=6.5*inch, height=4.5*inch))
                            story.append(Spacer(1, 0.3*inch))
                        except Exception as e:
                            logger.warning(f"Could not add chart {chart_path}: {str(e)}")
            
            story.append(PageBreak())
            story.append(Paragraph("Appendix", heading2_style))
            story.append(Paragraph("Report Metadata", heading3_style))
            story.append(Paragraph(f"Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", normal_style))
            story.append(Paragraph(f"Total Sections: {len(sections)}", normal_style))
            story.append(Paragraph(f"Total Charts: {len(charts) if charts else 0}", normal_style))
            story.append(Paragraph(f"Total Tables: {len(tables) if tables else 0}", normal_style))
            
            doc.build(story)
            logger.info(f"Generated comprehensive PDF report: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Error generating PDF: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            raise

    def generate_powerpoint_report(
        self,
        title: str,
        subtitle: str,
        slides_content: List[Dict[str, Any]],
        output_filename: str = "report.pptx"
    ) -> str:
        """
        Generate a PowerPoint report.
        
        Args:
            title: Presentation title
            subtitle: Presentation subtitle
            slides_content: List of slide dictionaries with content
            output_filename: Output PPTX filename
            
        Returns:
            Path to generated PPTX
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            DARK_BLUE = RGBColor(46, 134, 171)
            LIGHT_GRAY = RGBColor(240, 240, 240)
            WHITE = RGBColor(255, 255, 255)
            DARK_TEXT = RGBColor(50, 50, 50)
            
            title_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(title_slide_layout)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = DARK_BLUE
            
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(9)
            height = Inches(1.5)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(54)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = WHITE
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            subtitle_box = slide.shapes.add_textbox(left, Inches(4.2), width, Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(28)
            subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            date_box = slide.shapes.add_textbox(left, Inches(6.5), width, Inches(0.5))
            date_frame = date_box.text_frame
            date_frame.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
            date_frame.paragraphs[0].font.size = Pt(14)
            date_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            for slide_data in slides_content:
                self._add_content_slide_with_overflow(prs, slide_data, DARK_BLUE, LIGHT_GRAY, DARK_TEXT)
            
            output_path = self.output_dir / output_filename
            prs.save(str(output_path))
            logger.info(f"Generated PowerPoint report: {output_path}")
            return str(output_path)
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint: {str(e)}")
            raise

    def _add_content_slide_with_overflow(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        dark_blue: RGBColor,
        light_gray: RGBColor,
        dark_text: RGBColor
    ) -> None:
        """
        Add content slide(s) with automatic overflow handling.
        Splits content across multiple slides if needed.
        """
        content = slide_data.get('content', '')
        title_text = slide_data.get('title', 'Slide')
        
        MAX_ITEMS_PER_SLIDE = 7
        MAX_CHARACTERS_PER_SLIDE = 600
        MAX_LINES_PER_SLIDE = 12
        
        if isinstance(content, list):
            total_items = len(content)
            slide_num = 1
            
            for start_idx in range(0, total_items, MAX_ITEMS_PER_SLIDE):
                end_idx = min(start_idx + MAX_ITEMS_PER_SLIDE, total_items)
                items_for_slide = content[start_idx:end_idx]
                
                slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                slide_image = slide_data.get('image') if slide_num == 1 else None
                self._add_content_slide(
                    prs,
                    {
                        'title': slide_title,
                        'content': items_for_slide,
                        'image': slide_image
                    },
                    dark_blue,
                    light_gray,
                    dark_text
                )
                slide_num += 1
        else:
            content_str = clean_text_for_presentation(str(content))
            
            if len(content_str) <= MAX_CHARACTERS_PER_SLIDE:
                self._add_content_slide(prs, slide_data, dark_blue, light_gray, dark_text)
            else:
                paragraphs = content_str.split('\n\n')
                current_slide_content = []
                current_length = 0
                slide_num = 1
                
                for para in paragraphs:
                    para = para.strip()
                    if not para:
                        continue
                    
                    para_length = len(para)
                    
                    if current_length + para_length > MAX_CHARACTERS_PER_SLIDE and current_slide_content:
                        slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                        self._add_content_slide(
                            prs,
                            {
                                'title': slide_title,
                                'content': '\n\n'.join(current_slide_content),
                                'image': slide_data.get('image') if slide_num == 1 else None
                            },
                            dark_blue,
                            light_gray,
                            dark_text
                        )
                        current_slide_content = []
                        current_length = 0
                        slide_num += 1
                    
                    if para_length > MAX_CHARACTERS_PER_SLIDE:
                        sentences = para.split('. ')
                        for sentence in sentences:
                            sentence = sentence.strip()
                            if not sentence:
                                continue
                            if not sentence.endswith('.'):
                                sentence += '.'
                            
                            sentence_length = len(sentence)
                            if current_length + sentence_length > MAX_CHARACTERS_PER_SLIDE and current_slide_content:
                                slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                                slide_image = slide_data.get('image') if slide_num == 1 else None
                                self._add_content_slide(
                                    prs,
                                    {
                                        'title': slide_title,
                                        'content': '\n\n'.join(current_slide_content),
                                        'image': slide_image
                                    },
                                    dark_blue,
                                    light_gray,
                                    dark_text
                                )
                                current_slide_content = []
                                current_length = 0
                                slide_num += 1
                            
                            if sentence_length > MAX_CHARACTERS_PER_SLIDE:
                                words = sentence.split(' ')
                                temp_sentence = []
                                temp_length = 0
                                
                                for word in words:
                                    word_length = len(word) + 1
                                    if current_length + temp_length + word_length > MAX_CHARACTERS_PER_SLIDE and (temp_sentence or current_slide_content):
                                        if temp_sentence:
                                            current_slide_content.append(' '.join(temp_sentence))
                                            current_length += temp_length + 2
                                            temp_sentence = []
                                            temp_length = 0
                                        
                                        if current_slide_content:
                                            slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                                            slide_image = slide_data.get('image') if slide_num == 1 else None
                                            self._add_content_slide(
                                                prs,
                                                {
                                                    'title': slide_title,
                                                    'content': '\n\n'.join(current_slide_content),
                                                    'image': slide_image
                                                },
                                                dark_blue,
                                                light_gray,
                                                dark_text
                                            )
                                            current_slide_content = []
                                            current_length = 0
                                            slide_num += 1
                                    
                                    temp_sentence.append(word)
                                    temp_length += word_length
                                
                                if temp_sentence:
                                    if current_length + temp_length > MAX_CHARACTERS_PER_SLIDE and current_slide_content:
                                        slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                                        slide_image = slide_data.get('image') if slide_num == 1 else None
                                        self._add_content_slide(
                                            prs,
                                            {
                                                'title': slide_title,
                                                'content': '\n\n'.join(current_slide_content),
                                                'image': slide_image
                                            },
                                            dark_blue,
                                            light_gray,
                                            dark_text
                                        )
                                        current_slide_content = []
                                        current_length = 0
                                        slide_num += 1
                                    
                                    current_slide_content.append(' '.join(temp_sentence))
                                    current_length += temp_length + 2
                            else:
                                current_slide_content.append(sentence)
                                current_length += sentence_length + 2
                    else:
                        current_slide_content.append(para)
                        current_length += para_length + 2
                
                if current_slide_content:
                    slide_title = title_text if slide_num == 1 else f"{title_text} (continued)"
                    self._add_content_slide(
                        prs,
                        {
                            'title': slide_title,
                            'content': '\n\n'.join(current_slide_content),
                            'image': slide_data.get('image') if slide_num == 1 else None
                        },
                        dark_blue,
                        light_gray,
                        dark_text
                    )

    def _add_content_slide(
        self,
        prs: Presentation,
        slide_data: Dict[str, Any],
        dark_blue: RGBColor,
        light_gray: RGBColor,
        dark_text: RGBColor
    ) -> None:
        """Add a single content slide to presentation with improved formatting."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = dark_blue
        header_shape.line.color.rgb = dark_blue
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.9))
        title_frame = title_box.text_frame
        title_text = slide_data.get('title', 'Slide')
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        title_frame.margin_bottom = Inches(0.1)
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = None
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        
        content = slide_data.get('content', '')
        
        if isinstance(content, list):
            for idx, item in enumerate(content):
                if idx > 0:
                    text_frame.add_paragraph()
                
                item_clean = clean_text_for_presentation(str(item))
                if not item_clean:
                    continue
                
                p = text_frame.paragraphs[len(text_frame.paragraphs) - 1]
                item_clean = item_clean.lstrip('•-*0123456789.) ').strip()
                p.text = f"• {item_clean}"
                p.font.size = Pt(20)
                p.font.color.rgb = dark_text
                p.level = 0
                p.space_after = Pt(12)
                p.line_spacing = 1.2
        else:
            content_str = clean_text_for_presentation(str(content))
            paragraphs = content_str.split('\n\n')
            
            for para_idx, para in enumerate(paragraphs):
                if para_idx > 0:
                    text_frame.add_paragraph()
                
                para = para.strip()
                if not para:
                    continue
                
                if para.startswith('•') or para.startswith('-') or para.startswith('*'):
                    para = para.lstrip('•-* ').strip()
                    p = text_frame.paragraphs[len(text_frame.paragraphs) - 1]
                    p.text = f"• {para}"
                    p.font.size = Pt(20)
                    p.font.color.rgb = dark_text
                    p.level = 0
                    p.space_after = Pt(12)
                else:
                    lines = para.split('\n')
                    for line_idx, line in enumerate(lines):
                        if line_idx > 0:
                            text_frame.add_paragraph()
                        
                        line = line.strip()
                        if not line:
                            continue
                        
                        line = clean_text_for_presentation(line)
                        p = text_frame.paragraphs[len(text_frame.paragraphs) - 1]
                        p.text = line
                        p.font.size = Pt(20)
                        p.font.color.rgb = dark_text
                        p.level = 0
                        p.space_after = Pt(10) if line_idx < len(lines) - 1 else Pt(12)
                        p.line_spacing = 1.3
        
        if 'image' in slide_data and slide_data.get('image') and os.path.exists(slide_data['image']):
            try:
                content_length = len(str(content)) if isinstance(content, str) else sum(len(str(item)) for item in content) if isinstance(content, list) else 0
                if content_length < 100:
                    slide.shapes.add_picture(
                        slide_data['image'],
                        Inches(0.5),
                        Inches(1.5),
                        width=Inches(9),
                        height=Inches(5.5)
                    )
                else:
                    slide.shapes.add_picture(
                        slide_data['image'],
                        Inches(5.5),
                        Inches(3),
                        width=Inches(4),
                        height=Inches(3)
                    )
            except Exception as e:
                logger.warning(f"Could not add image: {str(e)}")

    def cleanup_temp_files(self) -> None:
        """Clean up temporary chart files."""
        try:
            for chart_path in self.charts:
                if os.path.exists(chart_path):
                    os.remove(chart_path)
            logger.info("Cleaned up temporary files")
        except Exception as e:
            logger.warning(f"Error cleaning up temporary files: {str(e)}")
